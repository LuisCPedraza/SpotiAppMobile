from pptx import Presentation
import sys

path = r"c:\Users\carlo\Proyectos\SopotiAppMobile\docs\spotiApp.pptx"
try:
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        title = title_shape.text if title_shape else "No Title"
        # Extract remaining slides
        if i >= 10:
            print(f"--- Slide {i} ---")
            print(f"Title: {title}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            print(f"  - {text}")
            print()
except Exception as e:
    print(f"Error: {e}")
